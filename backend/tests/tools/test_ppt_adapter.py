from uuid import uuid4

import pptx
import pytest
from shared.contracts.task import Task, TaskCategory

from app.tools.ppt.adapter import PPTToolAdapter


@pytest.mark.asyncio
async def test_ppt_tool_adapter_5_slides(tmp_path):
    adapter = PPTToolAdapter()

    out_file = str(tmp_path / "custom_5_slides.pptx")
    task = Task(
        workflow_id=uuid4(),
        task_name="Generate EV Presentation",
        description="Create presentation deck for EV Market Analysis",
        required_tool="ppt_tool",
        category=TaskCategory.PPT_GENERATION,
        expected_output="5-slide presentation deck",
        inputs={
            "topic": "EV Market Analysis 2026",
            "output_path": out_file,
        },
    )

    result = await adapter.execute(task)

    assert result.success is True
    assert len(result.artifacts) == 1
    art = result.artifacts[0]
    assert art.name == "custom_5_slides.pptx"

    # Verify python-pptx can open the created presentation and contains 5 slides
    prs = pptx.Presentation(out_file)
    assert len(prs.slides) == 5
    assert prs.slides[0].shapes.title.text == "EV Market Analysis 2026"
