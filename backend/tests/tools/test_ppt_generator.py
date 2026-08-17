import os
from uuid import uuid4

import pptx
import pytest
from pydantic import ValidationError
from shared.contracts.permission import (
    PermissionRequest,
    PermissionStatus,
)
from shared.contracts.tool import ToolHealth, ToolState

from app.core.exceptions import PermissionDeniedException
from app.core.permissions.manager import PermissionManager
from app.schemas.ppt import PresentationSchema, SlideContent, SlideType
from app.tools.ppt import PPTGenerator, ppt_tool_metadata
from app.tools.registry import ToolRegistry


def test_tool_metadata_and_registration():
    """Verify tool metadata is defined correctly
    and can be registered in ToolRegistry.
    """
    registry = ToolRegistry()
    registry.register(ppt_tool_metadata)

    tool = registry.get("ppt_tool")
    assert tool is not None
    assert tool.name == "ppt_tool"
    assert tool.status == ToolState.READY
    assert tool.health == ToolHealth.HEALTHY
    assert "python-pptx" in tool.dependencies
    assert "FILE_SYSTEM" in tool.required_permissions


def test_single_slide_generation(tmp_path):
    """Test generating a single title slide."""
    output_path = str(tmp_path / "single.pptx")
    workflow_id = uuid4()

    schema = PresentationSchema(
        title="Intro to AetherPhoenix",
        subtitle="V1 presentation",
        slides=[
            SlideContent(
                slide_type=SlideType.TITLE,
                title="AetherPhoenix Title Slide",
                subtitle="The sub-title text",
                speaker_notes="Welcome to AetherPhoenix!",
            )
        ],
    )

    generator = PPTGenerator()
    result = generator.generate(schema, output_path, workflow_id)

    assert result.file_path == os.path.abspath(output_path)
    assert result.slide_count == 1
    assert result.file_size > 0

    # Read back and inspect
    prs = pptx.Presentation(output_path)
    assert len(prs.slides) == 1

    title_slide = prs.slides[0]
    assert title_slide.shapes.title.text == "AetherPhoenix Title Slide"
    assert title_slide.placeholders[1].text == "The sub-title text"
    assert title_slide.notes_slide.notes_text_frame.text == "Welcome to AetherPhoenix!"


def test_multi_slide_generation(tmp_path):
    """Test generating a title slide + multiple content slides."""
    output_path = str(tmp_path / "multi.pptx")
    workflow_id = uuid4()

    schema = PresentationSchema(
        title="System Architecture",
        slides=[
            SlideContent(
                slide_type=SlideType.TITLE,
                title="System Overview",
                subtitle="AetherPhoenix architecture",
            ),
            SlideContent(
                slide_type=SlideType.CONTENT,
                title="Runtime Layer",
                bullets=[
                    "Central runtime engine handles tasks",
                    "Compiles plans from Planner Agent",
                    "Dispatches commands to Sandbox",
                ],
                speaker_notes="This is the main loop.",
            ),
            SlideContent(
                slide_type=SlideType.CONTENT,
                title="Healing Layer",
                bullets=[
                    "Self-heals runtime execution exceptions",
                    "Maintains state rollback history",
                ],
            ),
        ],
    )

    generator = PPTGenerator()
    result = generator.generate(schema, output_path, workflow_id)

    assert result.slide_count == 3

    # Read back and inspect
    prs = pptx.Presentation(output_path)
    assert len(prs.slides) == 3

    # Check slide 2
    slide2 = prs.slides[1]
    assert slide2.shapes.title.text == "Runtime Layer"
    # Check bullets text
    bullets_tf = slide2.placeholders[1].text_frame
    paragraphs = [p.text for p in bullets_tf.paragraphs]
    assert len(paragraphs) == 3
    assert paragraphs[0] == "Central runtime engine handles tasks"
    assert paragraphs[1] == "Compiles plans from Planner Agent"
    assert paragraphs[2] == "Dispatches commands to Sandbox"
    assert slide2.notes_slide.notes_text_frame.text == "This is the main loop."


def test_empty_bullets_content_slide(tmp_path):
    """Verify that a content slide can have no bullets (or empty bullets list)."""
    output_path = str(tmp_path / "empty_bullets.pptx")
    workflow_id = uuid4()

    schema = PresentationSchema(
        title="Empty Slide Title",
        slides=[
            SlideContent(
                slide_type=SlideType.CONTENT,
                title="Empty Content Slide",
                bullets=[],
            )
        ],
    )

    generator = PPTGenerator()
    result = generator.generate(schema, output_path, workflow_id)

    assert result.slide_count == 1

    prs = pptx.Presentation(output_path)
    slide = prs.slides[0]
    assert slide.shapes.title.text == "Empty Content Slide"
    # The text frame should not be populated with text if bullets are empty
    # Placeholders length can be checked
    assert len(slide.placeholders) >= 1


def test_validation_errors():
    """Verify that validation errors are raised
    for invalid presentation schema inputs.
    """
    # Min length check on presentation title
    with pytest.raises(ValidationError):
        PresentationSchema(title="", slides=[])

    # Min length check on slide title
    with pytest.raises(ValidationError):
        PresentationSchema(
            title="Valid Title",
            slides=[SlideContent(slide_type=SlideType.CONTENT, title="")],
        )


def test_permission_denied_exception(tmp_path):
    """Verify that if PermissionManager rejects the FILE_SYSTEM permission,
    we raise PermissionDeniedException.
    """
    output_path = str(tmp_path / "denied.pptx")
    workflow_id = uuid4()

    schema = PresentationSchema(
        title="Security Test",
        slides=[SlideContent(slide_type=SlideType.TITLE, title="Security Check")],
    )

    # Custom PermissionManager that rejects all requests
    class RejectingPermissionManager(PermissionManager):
        def request_permission(self, request: PermissionRequest) -> PermissionRequest:
            request.status = PermissionStatus.REJECTED
            return request

    pm = RejectingPermissionManager()
    generator = PPTGenerator(permission_manager=pm)

    with pytest.raises(PermissionDeniedException) as excinfo:
        generator.generate(schema, output_path, workflow_id)

    assert "FILE_SYSTEM permission denied" in str(excinfo.value)
    assert excinfo.value.code == "PERMISSION_DENIED"
    assert not os.path.exists(output_path)
