import os
from datetime import datetime, timezone
from uuid import UUID

import pptx
from pptx.dml.color import RGBColor
from pptx.util import Pt
from shared.contracts.permission import (
    PermissionRequest,
    PermissionStatus,
    PermissionType,
    RiskLevel,
)

from app.core.exceptions import PermissionDeniedException, ToolExecutionException
from app.core.logging import get_logger
from app.core.permissions.manager import PermissionManager
from app.schemas.ppt import PPTGenerationResult, PresentationSchema, SlideType

logger = get_logger(__name__)


class PPTGenerator:
    """
    Worker Agent tool for compiling structured presentation data
    into a PowerPoint (.pptx) file.
    Includes basic formatting, file validation, and permission checks.
    """

    def __init__(self, permission_manager: PermissionManager | None = None) -> None:
        self.permission_manager = permission_manager or PermissionManager()

    def generate(
        self, presentation: PresentationSchema, output_path: str, workflow_id: UUID
    ) -> PPTGenerationResult:
        """
        Generates a PowerPoint presentation at the specified output path.

        Parameters
        ----------
        presentation:
            The structured content and slide data.
        output_path:
            Target file path on the filesystem.
        workflow_id:
            The unique identifier of the active workflow (for permissions).

        Returns
        -------
        PPTGenerationResult:
            Metadata about the generated PowerPoint file.
        """
        logger.info(
            f"Initiating PPT generation: Workflow={workflow_id}, Output={output_path}"
        )

        # 1. Permission Check
        perm_req = PermissionRequest(
            workflow_id=workflow_id,
            permission_type=PermissionType.FILE_SYSTEM,
            reason="Generate PowerPoint presentation file",
            risk_level=RiskLevel.SAFE,
        )
        granted_req = self.permission_manager.request_permission(perm_req)

        if granted_req.status != PermissionStatus.GRANTED:
            logger.error("FILE_SYSTEM write permission rejected.")
            raise PermissionDeniedException(
                message="Cannot write PowerPoint file: FILE_SYSTEM permission denied.",
                details={"permission_type": PermissionType.FILE_SYSTEM},
            )

        # 2. PowerPoint Generation
        try:
            prs = pptx.Presentation()

            # Set slide width and height to default 16:9 widescreen layout
            prs.slide_width = Pt(960)  # 13.33 inches
            prs.slide_height = Pt(540)  # 7.5 inches

            # Iterate over the slide structures
            for idx, slide_data in enumerate(presentation.slides):
                logger.debug(
                    f"Compiling slide {idx + 1}/{len(presentation.slides)}: "
                    f"Title='{slide_data.title}'"
                )

                if slide_data.slide_type == SlideType.TITLE:
                    # Layout 0 is the default Title Slide layout
                    slide = prs.slides.add_slide(prs.slide_layouts[0])
                    self._populate_title_slide(slide, slide_data)
                else:
                    # Layout 1 is the default Title and Content layout
                    slide = prs.slides.add_slide(prs.slide_layouts[1])
                    self._populate_content_slide(slide, slide_data)

                # Append speaker notes if provided
                if slide_data.speaker_notes:
                    notes_slide = slide.notes_slide
                    notes_slide.notes_text_frame.text = slide_data.speaker_notes

            # Ensure directory exists before saving
            dir_name = os.path.dirname(output_path)
            if dir_name and not os.path.exists(dir_name):
                os.makedirs(dir_name, exist_ok=True)

            prs.save(output_path)
            logger.info(f"Presentation saved successfully to {output_path}")

        except Exception as e:
            logger.exception("Failed to compile PowerPoint deck.")
            raise ToolExecutionException(
                message=f"PowerPoint generation failed: {str(e)}",
                details={"output_path": output_path},
            )

        # 3. File Validation
        self._validate_output_file(output_path, len(presentation.slides))

        # 4. Generate Metadata
        file_size = os.path.getsize(output_path)
        return PPTGenerationResult(
            file_path=os.path.abspath(output_path),
            file_size=file_size,
            slide_count=len(presentation.slides),
            generated_at=datetime.now(timezone.utc),
        )

    def _populate_title_slide(self, slide, slide_data) -> None:
        """Sets text and styles for a title slide layout."""
        title_shape = slide.shapes.title
        subtitle_shape = slide.placeholders[1] if len(slide.placeholders) > 1 else None

        if title_shape:
            title_shape.text = slide_data.title
            self._apply_text_formatting(
                title_shape, font_size=Pt(44), bold=True, color=RGBColor(26, 54, 93)
            )

        if subtitle_shape and slide_data.subtitle:
            subtitle_shape.text = slide_data.subtitle
            self._apply_text_formatting(
                subtitle_shape,
                font_size=Pt(24),
                bold=False,
                color=RGBColor(74, 85, 104),
            )

    def _populate_content_slide(self, slide, slide_data) -> None:
        """Sets text, bullet points, and styles for content slide layout."""
        title_shape = slide.shapes.title
        content_shape = slide.placeholders[1] if len(slide.placeholders) > 1 else None

        if title_shape:
            title_shape.text = slide_data.title
            self._apply_text_formatting(
                title_shape, font_size=Pt(36), bold=True, color=RGBColor(26, 54, 93)
            )

        if content_shape and slide_data.bullets:
            tf = content_shape.text_frame
            # Reuse the first paragraph created by default
            first_p = tf.paragraphs[0]
            first_p.text = slide_data.bullets[0]
            self._apply_paragraph_formatting(first_p, font_size=Pt(18))

            for bullet in slide_data.bullets[1:]:
                p = tf.add_paragraph()
                p.text = bullet
                self._apply_paragraph_formatting(p, font_size=Pt(18))

    def _apply_text_formatting(self, shape, font_size, bold, color) -> None:
        """Applies consistent font face, size, and coloring to a shape's text frame."""
        if not shape.has_text_frame:
            return
        tf = shape.text_frame
        for p in tf.paragraphs:
            self._apply_paragraph_formatting(p, font_size, bold, color)

    def _apply_paragraph_formatting(
        self, paragraph, font_size, bold=False, color=None
    ) -> None:
        """Helper to style runs within a single paragraph."""
        paragraph.font.name = "Arial"
        paragraph.font.size = font_size
        paragraph.font.bold = bold
        if color:
            paragraph.font.color.rgb = color

    def _validate_output_file(self, file_path: str, expected_slides: int) -> None:
        """Checks the output path to verify the file is a valid PowerPoint deck."""
        logger.debug(f"Verifying generated PowerPoint deck at: {file_path}")

        if not os.path.exists(file_path):
            raise ToolExecutionException(
                message=f"Validation failed: File was not created at {file_path}."
            )

        if os.path.getsize(file_path) == 0:
            raise ToolExecutionException(
                message=f"Validation failed: File created at {file_path} is empty."
            )

        try:
            prs = pptx.Presentation(file_path)
            slide_count = len(prs.slides)
            if slide_count != expected_slides:
                raise ToolExecutionException(
                    message=(
                        f"Validation failed: Slide count mismatch. "
                        f"Expected {expected_slides}, got {slide_count}."
                    )
                )
            logger.info("PowerPoint verification check succeeded.")
        except Exception as e:
            logger.exception("Failed to parse output PPTX file.")
            raise ToolExecutionException(
                message=(
                    f"Validation failed: Generated file is corrupted or invalid. "
                    f"Details: {str(e)}"
                )
            )
