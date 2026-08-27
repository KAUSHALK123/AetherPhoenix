import hashlib
import inspect
import json
import logging
import time
from datetime import datetime, timezone
from pathlib import Path
from typing import Any
from uuid import UUID

from PIL import Image
from shared.contracts.artifact import Artifact, ArtifactType
from shared.contracts.document import DocumentFormat, StructuredDocumentInput
from shared.contracts.export import ExportFormat, ExportRequest, ExportResult
from shared.contracts.pdf import (
    HeadingElement,
    ParagraphElement,
    PDFDocumentInput,
    PDFElement,
)
from shared.contracts.permission import PermissionType

from app.core.config import get_config
from app.core.exceptions import PermissionDeniedException
from app.services.artifact_storage import (
    ArtifactStorageService,
    get_artifact_storage_service,
)
from app.tools.document.generator import DocumentGenerator
from app.tools.pdf.generator import PDFGenerator

logger = logging.getLogger(__name__)


class ExportError(Exception):
    """Raised when artifact conversion, export, or validation fails."""

    pass


class ExportEngine:
    """
    Unified Export Engine responsible for converting task artifacts
    across supported formats (PPTX, PDF, Markdown, DOCX, HTML, Images, CSV, JSON, TXT),
    preserving metadata, validating export health, and registering exported artifacts.
    """

    FORMAT_MAP: dict[ExportFormat, ArtifactType] = {
        ExportFormat.PPTX: ArtifactType.PPT,
        ExportFormat.PDF: ArtifactType.PDF,
        ExportFormat.MARKDOWN: ArtifactType.REPORTS,
        ExportFormat.DOCX: ArtifactType.REPORTS,
        ExportFormat.HTML: ArtifactType.REPORTS,
        ExportFormat.TXT: ArtifactType.REPORTS,
        ExportFormat.JSON: ArtifactType.DATA,
        ExportFormat.CSV: ArtifactType.DATA,
        ExportFormat.PNG: ArtifactType.IMAGES,
        ExportFormat.JPEG: ArtifactType.IMAGES,
        ExportFormat.WEBP: ArtifactType.IMAGES,
    }

    def __init__(
        self,
        permission_manager: Any | None = None,
        artifact_storage_service: ArtifactStorageService | None = None,
    ) -> None:
        self.permission_manager = permission_manager
        self.artifact_storage_service = (
            artifact_storage_service or get_artifact_storage_service()
        )
        self.pdf_generator = PDFGenerator(permission_manager=permission_manager)
        self.doc_generator = DocumentGenerator(permission_manager=permission_manager)

    async def _check_permission(
        self,
        workflow_id: UUID | None = None,
        task_id: UUID | None = None,
    ) -> None:
        """Enforces permission validation for filesystem export operations."""
        if not self.permission_manager:
            return

        perm_to_check = PermissionType.FILE_SYSTEM
        wf_str = str(workflow_id) if workflow_id else "global"
        action = "Export workflow artifact"

        try:
            if hasattr(self.permission_manager, "check_permission"):
                check_fn = self.permission_manager.check_permission
                if inspect.iscoroutinefunction(check_fn):
                    is_granted = await check_fn(
                        action=action,
                        permission_type=perm_to_check,
                        workflow_id=wf_str,
                        context={"task_id": str(task_id) if task_id else None},
                    )
                else:
                    res = check_fn(
                        action=action,
                        permission_type=perm_to_check,
                        workflow_id=wf_str,
                    )
                    if hasattr(res, "__await__"):
                        is_granted = await res
                    else:
                        is_granted = bool(res)

                if not is_granted:
                    # Fallback check with FILE_WRITE or FILE_SYSTEM_WRITE
                    for alt_perm in [
                        PermissionType.FILE_WRITE,
                        PermissionType.FILE_SYSTEM_WRITE,
                    ]:
                        alt_res = self.permission_manager.check_permission(
                            action=action,
                            permission_type=alt_perm,
                            workflow_id=wf_str,
                        )
                        if hasattr(alt_res, "__await__"):
                            is_granted = await alt_res
                        else:
                            is_granted = bool(alt_res)
                        if is_granted:
                            break

                if not is_granted:
                    raise PermissionDeniedException(
                        f"Permission denied: Missing '{perm_to_check.value}' "
                        "for export."
                    )
        except PermissionDeniedException:
            raise
        except Exception as e:
            logger.warning(f"Permission check encountered error during export: {e}")
            raise PermissionDeniedException(
                f"Permission check failed for '{perm_to_check.value}': {str(e)}"
            )

    async def export(self, request: ExportRequest) -> ExportResult:
        """
        Executes unified export for a source artifact or task result payload.

        Args:
            request: Export parameters.

        Returns:
            ExportResult metadata contract.
        """
        start_time = time.time()
        logger.info(
            f"Initiating artifact export: Workflow={request.workflow_id}, "
            f"TargetFormat={request.target_format.value}"
        )

        # 1. Permission Enforcement
        await self._check_permission(
            workflow_id=request.workflow_id, task_id=request.task_id
        )

        # 2. Resolve Source Content and Metadata
        (
            source_content,
            source_path,
            source_artifact_id,
            source_metadata,
        ) = await self._resolve_source(request)

        # 3. Determine Output File Path
        output_filepath = self._determine_output_path(request)

        # 4. Perform Format Conversion
        try:
            target_path = await self._convert_format(
                source_content=source_content,
                source_path=source_path,
                target_format=request.target_format,
                output_path=output_filepath,
                request=request,
                source_metadata=source_metadata,
            )
        except Exception as conv_err:
            logger.error(
                f"Format conversion to {request.target_format.value} "
                f"failed: {conv_err}",
                exc_info=True,
            )
            raise ExportError(
                f"Conversion to {request.target_format.value} failed: {str(conv_err)}"
            ) from conv_err

        # 5. File Validation (Verify file exists, non-empty, and readable)
        self._validate_exported_file(target_path)

        # 6. Calculate File Size & SHA-256 Checksum
        file_bytes = target_path.read_bytes()
        file_size = len(file_bytes)
        checksum = hashlib.sha256(file_bytes).hexdigest()

        artifact_name = request.title or target_path.name
        art_type = self.FORMAT_MAP.get(request.target_format, ArtifactType.REPORTS)

        source_art_id_str = str(source_artifact_id) if source_artifact_id else None
        source_fp_str = str(source_path) if source_path else None

        # Merge metadata preserving source lineage
        merged_metadata = {
            **source_metadata,
            **request.metadata,
            "export_format": request.target_format.value,
            "source_artifact_id": source_art_id_str,
            "source_filepath": source_fp_str,
            "exported_at": datetime.now(timezone.utc).isoformat(),
        }

        # 7. Register Exported Artifact under the Workflow ID
        exported_artifact = Artifact(
            workflow_id=request.workflow_id,
            task_id=request.task_id,
            name=artifact_name,
            filepath=str(target_path.resolve()),
            artifact_type=art_type,
            size_bytes=file_size,
            checksum=checksum,
            metadata=merged_metadata,
        )

        saved_artifact = await self.artifact_storage_service.register_artifact(
            artifact=exported_artifact,
            content=file_bytes,
        )

        elapsed_ms = (time.time() - start_time) * 1000.0
        download_url = f"/api/v1/artifacts/{saved_artifact.artifact_id}/download"

        logger.info(
            f"Successfully exported artifact {saved_artifact.artifact_id} "
            f"('{saved_artifact.name}', {file_size} bytes, {elapsed_ms:.1f}ms) "
            f"under workflow {request.workflow_id}"
        )

        return ExportResult(
            artifact_id=saved_artifact.artifact_id,
            workflow_id=request.workflow_id,
            task_id=request.task_id,
            name=saved_artifact.name,
            filepath=str(target_path.resolve()),
            download_url=download_url,
            format=request.target_format,
            size_bytes=file_size,
            checksum=checksum,
            source_artifact_id=source_art_id_str,
            source_filepath=source_fp_str,
            status="SUCCESS",
            created_at=datetime.now(timezone.utc),
            execution_time_ms=elapsed_ms,
            metadata=merged_metadata,
        )

    async def _resolve_source(
        self, request: ExportRequest
    ) -> tuple[bytes | str | None, Path | None, UUID | str | None, dict[str, Any]]:
        """Resolves source artifact content, filepath, and metadata."""
        source_content: bytes | str | None = None
        source_path: Path | None = None
        source_artifact_id: UUID | str | None = request.source_artifact_id
        source_metadata: dict[str, Any] = {}

        # 1. Resolve from source_artifact_id
        if request.source_artifact_id:
            art = await self.artifact_storage_service.get_artifact(
                request.source_artifact_id
            )
            if art:
                source_metadata = art.metadata.copy() if art.metadata else {}
                if Path(art.filepath).exists():
                    source_path = Path(art.filepath)
                source_content = (
                    await self.artifact_storage_service.get_artifact_content(
                        request.source_artifact_id
                    )
                )

        # 2. Resolve from source_filepath
        if not source_content and not source_path and request.source_filepath:
            sp = Path(request.source_filepath)
            if sp.exists():
                source_path = sp
                try:
                    source_content = sp.read_bytes()
                except Exception as e:
                    logger.warning(
                        f"Could not read source filepath '{sp}': {e}. "
                        "Using path reference."
                    )

        # 3. Resolve from request metadata / inputs payload
        if not source_content and not source_path:
            meta_content = (
                request.metadata.get("content")
                or request.metadata.get("text")
                or request.metadata.get("raw_content")
            )
            if meta_content:
                source_content = meta_content
            elif "filepath" in request.metadata:
                fp = Path(request.metadata["filepath"])
                if fp.exists():
                    source_path = fp
                    source_content = fp.read_bytes()

        if not source_content and not source_path:
            raise ExportError(
                "Source artifact or file not found. Provide a valid "
                "'source_artifact_id' or 'source_filepath'."
            )

        return source_content, source_path, source_artifact_id, source_metadata

    def _determine_output_path(self, request: ExportRequest) -> Path:
        """Determines target export destination path on local disk."""
        if request.output_path:
            out_p = Path(request.output_path).expanduser().resolve()
            out_p.parent.mkdir(parents=True, exist_ok=True)
            return out_p

        cfg = get_config()
        base_dir = Path(cfg.ARTIFACTS_DIR) / str(request.workflow_id)
        base_dir.mkdir(parents=True, exist_ok=True)

        ext = request.target_format.value
        title_slug = (
            request.title.lower().replace(" ", "_")
            if request.title
            else f"export_{int(time.time())}"
        )
        filename = f"{title_slug}.{ext}"
        return base_dir / filename

    async def _convert_format(
        self,
        source_content: bytes | str | None,
        source_path: Path | None,
        target_format: ExportFormat,
        output_path: Path,
        request: ExportRequest,
        source_metadata: dict[str, Any],
    ) -> Path:
        """Performs format conversion from source to target format."""
        source_ext = source_path.suffix.lower() if source_path else ""

        # A. Identity / Direct Copy (Same format)
        if source_path and (
            source_ext.replace(".", "") == target_format.value
            or (
                source_ext in (".jpg", ".jpeg")
                and target_format in (ExportFormat.JPEG, ExportFormat.PNG)
            )
        ):
            if source_path.resolve() != output_path.resolve():
                import shutil

                shutil.copy2(source_path, output_path)
            return output_path

        # B. Image to Image Conversion (PNG / JPEG / WEBP)
        if target_format in (ExportFormat.PNG, ExportFormat.JPEG, ExportFormat.WEBP):
            if source_path and source_ext in (
                ".png",
                ".jpg",
                ".jpeg",
                ".bmp",
                ".webp",
                ".tiff",
            ):
                with Image.open(source_path) as img:
                    save_format = target_format.value.upper()
                    if save_format == "JPEG" and img.mode in ("RGBA", "LA", "P"):
                        img = img.convert("RGB")
                    img.save(output_path, format=save_format)
                return output_path
            elif isinstance(source_content, bytes):
                import io

                with Image.open(io.BytesIO(source_content)) as img:
                    save_format = target_format.value.upper()
                    if save_format == "JPEG" and img.mode in ("RGBA", "LA", "P"):
                        img = img.convert("RGB")
                    img.save(output_path, format=save_format)
                return output_path

        # C. Conversion to PDF
        if target_format == ExportFormat.PDF:
            # 1. PPTX to PDF
            if source_ext == ".pptx" and source_path:
                return self._convert_pptx_to_pdf(source_path, output_path, request)
            # 2. Text / Markdown to PDF
            text_str = self._to_string(source_content, source_path)
            return self._convert_text_to_pdf(
                text_str, output_path, request, source_metadata
            )

        # D. Conversion to Markdown / Text / HTML / DOCX
        if target_format in (
            ExportFormat.MARKDOWN,
            ExportFormat.HTML,
            ExportFormat.DOCX,
            ExportFormat.TXT,
        ):
            text_str = ""
            if source_ext == ".pptx" and source_path:
                text_str = self._extract_text_from_pptx(source_path)
            else:
                text_str = self._to_string(source_content, source_path)

            doc_fmt_map = {
                ExportFormat.MARKDOWN: DocumentFormat.MARKDOWN,
                ExportFormat.TXT: DocumentFormat.TEXT,
                ExportFormat.HTML: DocumentFormat.HTML,
                ExportFormat.DOCX: DocumentFormat.MARKDOWN,  # DOCX rich document format
            }
            doc_fmt = doc_fmt_map.get(target_format, DocumentFormat.MARKDOWN)

            doc_title = (
                request.title or source_metadata.get("title") or "Export Document"
            )
            doc_input = StructuredDocumentInput(
                title=doc_title,
                content=text_str,
                format=doc_fmt,
                output_path=str(output_path),
                overwrite=True,
                workflow_id=request.workflow_id,
                task_id=request.task_id,
            )
            doc_res = self.doc_generator.generate(doc_input)
            return Path(doc_res.filepath)

        # E. Conversion to CSV / JSON
        if target_format in (ExportFormat.CSV, ExportFormat.JSON):
            text_str = self._to_string(source_content, source_path)
            if target_format == ExportFormat.JSON:
                try:
                    parsed_json = json.loads(text_str)
                    output_path.write_text(
                        json.dumps(parsed_json, indent=2), encoding="utf-8"
                    )
                except Exception:
                    json_data = {
                        "title": request.title or "Export Data",
                        "content": text_str,
                    }
                    output_path.write_text(
                        json.dumps(json_data, indent=2), encoding="utf-8"
                    )
                return output_path
            else:  # CSV
                doc_input = StructuredDocumentInput(
                    title=request.title or "Export CSV",
                    content=text_str,
                    format=DocumentFormat.CSV,
                    output_path=str(output_path),
                    overwrite=True,
                    workflow_id=request.workflow_id,
                    task_id=request.task_id,
                )
                doc_res = self.doc_generator.generate(doc_input)
                return Path(doc_res.filepath)

        # Default Fallback: Write text/bytes directly
        if isinstance(source_content, bytes):
            output_path.write_bytes(source_content)
        else:
            text_val = str(source_content or "")
            output_path.write_text(text_val, encoding="utf-8")

        return output_path

    def _convert_pptx_to_pdf(
        self, source_path: Path, output_path: Path, request: ExportRequest
    ) -> Path:
        """Converts a PowerPoint presentation (.pptx) file into a PDF document."""
        pdf_elements: list[PDFElement] = []
        doc_title = request.title or source_path.stem.replace("_", " ").title()

        try:
            import pptx

            prs = pptx.Presentation(str(source_path))
            for idx, slide in enumerate(prs.slides, 1):
                title_text = f"Slide {idx}"
                slide_bullets = []

                for shape in slide.shapes:
                    if shape.has_text_frame:
                        text = shape.text_frame.text.strip()
                        if text:
                            if shape == slide.shapes.title and text:
                                title_text = text
                            else:
                                lines = [
                                    line.strip()
                                    for line in text.splitlines()
                                    if line.strip()
                                ]
                                slide_bullets.extend(lines)

                pdf_elements.append(HeadingElement(text=title_text, level=2))
                for bullet in slide_bullets:
                    pdf_elements.append(ParagraphElement(text=f"• {bullet}"))

        except Exception as e:
            logger.warning(
                f"Failed to parse PPTX with python-pptx ({e}); "
                "fallback to text rendering."
            )
            pdf_elements.append(
                ParagraphElement(text=f"PowerPoint Presentation: {source_path.name}")
            )

        pdf_input = PDFDocumentInput(
            title=doc_title,
            output_path=str(output_path),
            elements=pdf_elements,
            workflow_id=request.workflow_id,
            task_id=request.task_id,
        )
        pdf_res = self.pdf_generator.generate(pdf_input)
        return Path(pdf_res.filepath)

    def _convert_text_to_pdf(
        self,
        text_content: str,
        output_path: Path,
        request: ExportRequest,
        source_metadata: dict[str, Any],
    ) -> Path:
        """Renders raw text or markdown lines into a formatted PDF document."""
        elements: list[PDFElement] = []
        doc_title = (
            request.title
            or source_metadata.get("title")
            or output_path.stem.replace("_", " ").title()
        )

        for line in text_content.splitlines():
            line_str = line.strip()
            if not line_str:
                continue
            if line_str.startswith("# "):
                elements.append(HeadingElement(text=line_str[2:], level=1))
            elif line_str.startswith("## "):
                elements.append(HeadingElement(text=line_str[3:], level=2))
            elif line_str.startswith("### "):
                elements.append(HeadingElement(text=line_str[4:], level=3))
            elif line_str.startswith("- ") or line_str.startswith("* "):
                elements.append(ParagraphElement(text=f"• {line_str[2:]}"))
            else:
                elements.append(ParagraphElement(text=line_str))

        if not elements:
            elements.append(ParagraphElement(text="No content available."))

        pdf_input = PDFDocumentInput(
            title=doc_title,
            output_path=str(output_path),
            elements=elements,
            workflow_id=request.workflow_id,
            task_id=request.task_id,
        )
        pdf_res = self.pdf_generator.generate(pdf_input)
        return Path(pdf_res.filepath)

    def _extract_text_from_pptx(self, source_path: Path) -> str:
        """Extracts plain text content from a PPTX file."""
        lines = []
        try:
            import pptx

            prs = pptx.Presentation(str(source_path))
            for idx, slide in enumerate(prs.slides, 1):
                lines.append(f"## Slide {idx}")
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        text = shape.text_frame.text.strip()
                        if text:
                            lines.append(text)
                lines.append("")
        except Exception as e:
            logger.warning(f"Could not extract text from PPTX {source_path}: {e}")
            lines.append(f"PowerPoint Presentation: {source_path.name}")
        return "\n".join(lines)

    def _to_string(
        self, content: bytes | str | None, filepath: Path | None = None
    ) -> str:
        """Utility to convert bytes/string/filepath into string text."""
        if isinstance(content, str):
            return content
        if isinstance(content, bytes):
            try:
                return content.decode("utf-8")
            except UnicodeDecodeError:
                return content.decode("latin-1", errors="ignore")
        if filepath and filepath.exists():
            try:
                return filepath.read_text(encoding="utf-8", errors="ignore")
            except Exception:
                return ""
        return ""

    def _validate_exported_file(self, file_path: Path) -> None:
        """
        Validates that the exported output file exists on disk,
        is non-empty, and is readable.
        """
        if not file_path.exists():
            raise ExportError(
                f"Export validation failed: File does not exist at '{file_path}'."
            )

        if not file_path.is_file():
            raise ExportError(
                f"Export validation failed: Path is not a file '{file_path}'."
            )

        size = file_path.stat().st_size
        if size == 0:
            raise ExportError(
                f"Export validation failed: Exported file '{file_path}' "
                "is empty (0 bytes)."
            )

        try:
            with open(file_path, "rb") as f:
                f.read(10)
        except Exception as e:
            raise ExportError(
                f"Export validation failed: File '{file_path}' is unreadable "
                f"({str(e)})."
            ) from e

        logger.debug(
            f"Export file validation check passed for '{file_path}' ({size} bytes)."
        )
